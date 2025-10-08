import os, io, zipfile
from dotenv import load_dotenv

from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from google import genai
from google.genai import types

load_dotenv(dotenv_path=r"B:\PROJECTS\CHATBOT_EDUCATION\.env.init")
gemini_key = os.getenv("gemini_generate")


client = genai.Client(api_key=gemini_key)

caption_image_prompt_template = """
You are an image captioning and content integration module. Below is the *CONTEXT* extracted from the file:
{context}

INSTRUCTIONS:
1. First, decide: is this image **relevant** to the CONTEXT or not?
   - "Relevant" = the image illustrates, directly presents, or provides data/visual support for the CONTEXT (charts, diagrams, models, product photos, etc.).
2. **If NOT relevant**, return exactly **<NO_CAPTION>** (this exact string, nothing else).
3. **If relevant**, return a single continuous text block that both:
   - (a) provides a detailed and objective description of the image, following these rules:
       * Identify the type of image (e.g., "bar chart", "line graph", "process diagram", "3D model", "product photo", etc.).
       * Describe the key visual elements connected to the CONTEXT: axes, axis labels (if readable), units, legend, main components/relationships, trends, outliers or highlights.
       * If the image contains readable text/labels: quote them verbatim in quotation marks.
       * Do NOT invent numbers or details not visible in the image. If you must approximate, prefix with **"Estimated:"**.
       * If something is unclear, explicitly say "not readable" or "not identifiable".
       * The length of the description should adapt to the richness of the image: more detail for complex visuals, concise for simple ones.
   - (b) integrates the image into the CONTEXT by elaborating, extending, or explaining how this visual supports, clarifies, or enriches the given content.
4. Language: use the **same language as the CONTEXT**; if CONTEXT is empty, default to English.
5. Tone: descriptive, neutral, objective.
6. **Output requirement**: return only one plain string — the integrated description+context continuation (if relevant) or **<NO_CAPTION>** (if not relevant). No JSON, no bullet points, no extra text.

FINAL RULE: Output exactly one plain text string.
"""

def call_vision_caption(image: Image.Image, context: str = "") -> str:
    """Call Gemini vision model to caption an image with context awareness."""
    if image.mode in ("RGBA", "LA", "P"):
        image = image.convert("RGB")

    buffered = io.BytesIO()
    image.save(buffered, format="JPEG")
    image_bytes = buffered.getvalue()
    image_part = types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg")

    # Inject context into prompt
    prompt = caption_image_prompt_template.format(context=context or "")

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=[image_part, prompt]
    )
    caption = response.text.strip()

    if caption == "<NO_CAPTION>" or not caption:
        return ""
    return caption


def process_pdf(file_path: str):
    doc = fitz.open(file_path)
    parts = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            parts.append(text)

        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b["type"] == 1 and "image" in b:
                img = Image.open(io.BytesIO(b["image"]))
                caption = call_vision_caption(img, context=text)
                if caption:
                    parts.append(caption)
    return parts


def process_pptx(file_path: str):
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        slide_text = " ".join([sh.text.strip() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()])
        if slide_text:
            parts.append(slide_text)

        for shape in slide.shapes:
            if shape.shape_type == 13 and shape.image:
                img = Image.open(io.BytesIO(shape.image.blob))
                caption = call_vision_caption(img, context=slide_text)
                if caption:
                    parts.append(caption)
    return parts


def process_docx(file_path: str):
    parts = []
    doc = Document(file_path)

    doc_text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
    if doc_text:
        parts.append(doc_text)

    with zipfile.ZipFile(file_path, "r") as z:
        for f in z.namelist():
            if f.startswith("word/media/"):
                data = z.read(f)
                img = Image.open(io.BytesIO(data))
                caption = call_vision_caption(img, context=doc_text)
                if caption:
                    parts.append(caption)
    return parts


def process_txt(file_path: str):
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read().strip()
    return [text] if text else []


def process_file(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return process_pdf(file_path)
    if ext == ".pptx":
        return process_pptx(file_path)
    if ext == ".docx":
        return process_docx(file_path)
    if ext == ".txt":
        return process_txt(file_path)
    raise ValueError(f"Unsupported extension: {ext}")

if __name__=='__main__':
    path = r'B:\PROJECTS\CHATBOT_EDUCATION\test\test.pdf'
    parts = process_file(path)
    print(parts)
    print(type(parts))
    print(parts[0])

    print(" ".join(parts))
    print(type(parts))